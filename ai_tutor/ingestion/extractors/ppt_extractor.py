"""Extract text and slide images from PowerPoint files."""

import io
import logging
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from scholera.config import settings

logger = logging.getLogger(__name__)


def extract_pptx(file_path: str, material_id: str) -> dict:
    """
    Returns same structure as pdf_extractor.extract_pdf.
    """
    path = Path(file_path)
    logger.info("Extracting PPTX: %s", path.name)

    prs = Presentation(str(path))
    images_dir = settings.images_dir / material_id
    images_dir.mkdir(parents=True, exist_ok=True)

    pages = []
    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1
        text_parts = []
        has_images = False
        has_equations = False
        slide_image_path = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        text_parts.append(line)
                        if _detect_equations(line):
                            has_equations = True

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_images = True
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes))
                    img_path = images_dir / f"slide_{page_num}_img.png"
                    img.save(str(img_path), "PNG")
                    slide_image_path = str(img_path)
                except Exception:
                    logger.debug("Could not extract image from slide %d", page_num)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    text_parts.append(row_text)

        text = "\n".join(text_parts)

        pages.append({
            "page_number": page_num,
            "text": text,
            "image_path": slide_image_path,
            "has_equations": has_equations,
            "has_images": has_images,
            "text_density": len(text),
        })

    return {
        "pages": pages,
        "full_markdown": "\n\n---\n\n".join(p["text"] for p in pages),
        "page_count": len(pages),
    }


def _detect_equations(text: str) -> bool:
    patterns = [r"\$.*?\$", r"\\frac", r"\\sum", r"\\int", r"=\s*\\", r"∑", r"∫", r"∂", r"→"]
    return any(re.search(p, text) for p in patterns)
